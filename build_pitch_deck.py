"""
OpenClaw 应用创新大赛 · 初赛答辩 PPT 生成脚本
南大数智安全官 NJU Guardian · 作者 · 2026-05
8 页结构，总讲稿 ~5 分 15 秒，南大紫主题
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 主题色 ----------
NJU_PURPLE = RGBColor(0x7B, 0x28, 0x7D)   # 南大紫
NJU_PURPLE_DARK = RGBColor(0x4A, 0x18, 0x52)
NJU_PURPLE_LIGHT = RGBColor(0xF3, 0xE8, 0xF5)
TEXT_DARK = RGBColor(0x2D, 0x2D, 0x2D)
TEXT_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xE2, 0x3A, 0x3A)
YELLOW = RGBColor(0xF5, 0xA5, 0x23)
GREEN = RGBColor(0x35, 0xA8, 0x5C)
BORDER_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY = RGBColor(0xFA, 0xF9, 0xFB)

# 中文字体：macOS 默认 PingFang SC；Windows 自动回落到微软雅黑
CN_FONT = "PingFang SC"
EN_FONT = "Helvetica"

# 16:9 画布：13.333 in × 7.5 in
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_line(shape, rgb=None, width=None):
    if rgb is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb
    if width:
        shape.line.width = width


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        set_fill(shp, fill)
    if line is None:
        shp.line.fill.background()
    else:
        set_line(shp, line, line_w)
    return shp


def add_round_rect(slide, left, top, width, height, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill is None:
        shp.fill.background()
    else:
        set_fill(shp, fill)
    if line is None:
        shp.line.fill.background()
    else:
        set_line(shp, line, line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, *,
             size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, font=CN_FONT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_text_in_shape(shape, text, *, size=14, bold=False, color=TEXT_DARK,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=CN_FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_page_header(slide, page_no, title, subtitle, total=8):
    # 顶部左侧紫色竖条 + 标题
    add_rect(slide, Inches(0.5), Inches(0.45), Inches(0.12), Inches(0.55), fill=NJU_PURPLE)
    add_text(slide, Inches(0.72), Inches(0.4), Inches(10), Inches(0.5),
             title, size=26, bold=True, color=NJU_PURPLE_DARK, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, Inches(0.72), Inches(0.92), Inches(11.5), Inches(0.35),
                 subtitle, size=13, color=TEXT_GRAY)
    # 右上页码
    add_text(slide, Inches(11.6), Inches(0.45), Inches(1.3), Inches(0.4),
             f"{page_no:02d} / {total:02d}", size=11, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)


def add_footer(slide, time_budget, agenda):
    # 底部细线
    add_rect(slide, Inches(0.5), Inches(7.05), Inches(12.3), Emu(8000),
             fill=NJU_PURPLE_LIGHT)
    add_text(slide, Inches(0.5), Inches(7.12), Inches(8.5), Inches(0.3),
             f"南大数智安全官 NJU Guardian · OpenClaw 应用创新大赛初赛答辩 · 2026.05",
             size=10, color=TEXT_GRAY)
    add_text(slide, Inches(9.0), Inches(7.12), Inches(3.8), Inches(0.3),
             f"⏱ 讲稿 {time_budget}  ·  {agenda}",
             size=10, color=NJU_PURPLE, align=PP_ALIGN.RIGHT, bold=True)


# ============================================================
#  PPT 主体
# ============================================================
def build_deck(output_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    # ------------------ 第 1 页 · 封面 ------------------
    slide = prs.slides.add_slide(blank_layout)
    # 整页紫色背景
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NJU_PURPLE_DARK)
    # 右下装饰圆
    add_rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(1.7), fill=NJU_PURPLE)
    # 顶部赛事标
    add_text(slide, Inches(0.7), Inches(0.5), Inches(12), Inches(0.4),
             "OpenClaw 应用创新大赛 · 初赛答辩", size=14, color=NJU_PURPLE_LIGHT, bold=True)
    add_text(slide, Inches(0.7), Inches(0.9), Inches(12), Inches(0.4),
             "OPENCLAW APPLICATION INNOVATION CONTEST · ROUND ONE",
             size=10, color=NJU_PURPLE_LIGHT, font=EN_FONT)

    # 主标题
    add_text(slide, Inches(0.7), Inches(2.2), Inches(12), Inches(1.2),
             "🛡 南大数智安全官", size=58, color=TEXT_LIGHT, bold=True)
    add_text(slide, Inches(0.7), Inches(3.4), Inches(12), Inches(0.7),
             "NJU Guardian", size=34, color=NJU_PURPLE_LIGHT, font=EN_FONT)
    # 副标题
    add_text(slide, Inches(0.7), Inches(4.2), Inches(12), Inches(0.55),
             "基于 OpenClaw 平台的校园场景化反诈 AI Agent",
             size=22, color=TEXT_LIGHT)
    # 横线
    add_rect(slide, Inches(0.7), Inches(4.9), Inches(2), Emu(15000), fill=NJU_PURPLE_LIGHT)
    # 关键指标条
    add_text(slide, Inches(0.7), Inches(5.05), Inches(12), Inches(0.5),
             "50 类知识库 · 双路召回 · 8/8 冒烟测试 · MIT + CC BY 4.0 开源",
             size=14, color=NJU_PURPLE_LIGHT)
    # 答辩人
    add_text(slide, Inches(0.7), Inches(6.05), Inches(8), Inches(0.4),
             "答辩人：作者（***REMOVED***）",
             size=15, color=TEXT_LIGHT, bold=True)
    add_text(slide, Inches(0.7), Inches(6.45), Inches(8), Inches(0.4),
             "南京大学 · 经济学（拔尖计划）· 个人参赛",
             size=12, color=NJU_PURPLE_LIGHT)
    # 右下角仓库
    add_text(slide, Inches(8), Inches(6.05), Inches(4.7), Inches(0.4),
             "github.com/hjjbh1314/nju-guardian",
             size=12, color=NJU_PURPLE_LIGHT, align=PP_ALIGN.RIGHT, font=EN_FONT)
    add_text(slide, Inches(8), Inches(6.45), Inches(4.7), Inches(0.4),
             "2026 · 南京", size=11, color=NJU_PURPLE_LIGHT, align=PP_ALIGN.RIGHT)

    # ------------------ 第 2 页 · 设计目的 ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, 2, "01  设计目的", "校园反诈的覆盖空白 · Why this project")

    # 三大痛点卡片（一行三个）
    pain_data = [
        ("📊 高发刑事犯罪", "电信网络诈骗已成为我国\n发案率最高的刑事犯罪之一",
         "大学生：社会经验不足、\n新型 App 使用频繁、属高危群体"),
        ("🎓 校园语境覆盖不足", "国家反诈中心 App 通用场景成熟",
         "但辅导员收班费、奖助、宿舍租赁、\n社团转账、老乡群高仿号——\n这些校园场景命中率明显偏低"),
        ("🚨 新生季高发四类", "冒充辅导员 / 学长索要班费",
         "新生兼职刷单（手册十大类首位）\n冒充客服+网贷+钓鱼链接\n高仿号+诱导安装陌生 App"),
    ]
    card_w = Inches(3.95)
    card_h = Inches(3.5)
    card_top = Inches(1.7)
    gap = Inches(0.15)
    left = Inches(0.5)
    for i, (title, lead, body) in enumerate(pain_data):
        x = left + (card_w + gap) * i
        card = add_round_rect(slide, x, card_top, card_w, card_h,
                              fill=BG_GRAY, line=NJU_PURPLE_LIGHT, line_w=Pt(0.75))
        card.adjustments[0] = 0.06
        add_text(slide, x + Inches(0.3), card_top + Inches(0.25), card_w - Inches(0.6), Inches(0.55),
                 title, size=18, bold=True, color=NJU_PURPLE)
        add_text(slide, x + Inches(0.3), card_top + Inches(0.95), card_w - Inches(0.6), Inches(0.8),
                 lead, size=13, color=TEXT_DARK, bold=True)
        add_text(slide, x + Inches(0.3), card_top + Inches(1.85), card_w - Inches(0.6), Inches(1.5),
                 body, size=12, color=TEXT_GRAY)

    # 底部结论条
    concl_top = Inches(5.55)
    bar = add_round_rect(slide, Inches(0.5), concl_top, Inches(12.3), Inches(1.0),
                         fill=NJU_PURPLE)
    bar.adjustments[0] = 0.15
    add_text_in_shape(bar,
                      "▶ 设计目的：补齐校园场景，做学生「粘进去就出结果」的工具——"
                      "短期补充国家反诈 App，长期嵌入南大官方 App 安全模块",
                      size=15, bold=True, color=TEXT_LIGHT,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, "00:40", "1/3 · 设计目的与思路")

    # ------------------ 第 3 页 · 核心思路 + 系统架构 ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, 3, "02  核心思路 · 系统架构",
                    "三模态输入 → OpenClaw 中间层 → 三段输出")

    # 三层架构图，纵向排列
    layer_left = Inches(0.5)
    layer_w = Inches(12.3)
    layer_h = Inches(1.5)

    # Layer 1 用户层
    top1 = Inches(1.55)
    l1 = add_round_rect(slide, layer_left, top1, layer_w, layer_h,
                        fill=NJU_PURPLE_LIGHT, line=NJU_PURPLE, line_w=Pt(1.0))
    l1.adjustments[0] = 0.05
    add_text(slide, layer_left + Inches(0.3), top1 + Inches(0.1), Inches(2.5), Inches(0.4),
             "① 用户层", size=15, bold=True, color=NJU_PURPLE_DARK)
    add_text(slide, layer_left + Inches(0.3), top1 + Inches(0.55), Inches(2.5), Inches(0.4),
             "微信小程序 / H5", size=11, color=TEXT_GRAY)
    # 三个输入入口
    inputs = [("📝 文本输入", "粘贴可疑短信、聊天记录"),
              ("🖼 截图上传", "OCR + 二维码识别"),
              ("🔗 链接电话", "可疑 URL / 手机号")]
    in_w = Inches(2.9)
    in_left = layer_left + Inches(3.2)
    for i, (k, v) in enumerate(inputs):
        x = in_left + (in_w + Inches(0.1)) * i
        b = add_round_rect(slide, x, top1 + Inches(0.2), in_w, Inches(1.1),
                           fill=TEXT_LIGHT, line=NJU_PURPLE, line_w=Pt(0.5))
        b.adjustments[0] = 0.1
        add_text(slide, x, top1 + Inches(0.3), in_w, Inches(0.5),
                 k, size=14, bold=True, color=NJU_PURPLE_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, x, top1 + Inches(0.8), in_w, Inches(0.4),
                 v, size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # 向下箭头
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(6.4), top1 + layer_h + Inches(0.05),
                                     Inches(0.5), Inches(0.4))
    set_fill(arrow1, NJU_PURPLE)
    arrow1.line.fill.background()

    # Layer 2 OpenClaw 中间层 —— 加重显示
    top2 = top1 + layer_h + Inches(0.55)
    layer_h2 = Inches(1.55)
    l2 = add_round_rect(slide, layer_left, top2, layer_w, layer_h2,
                        fill=NJU_PURPLE, line=NJU_PURPLE_DARK, line_w=Pt(1.5))
    l2.adjustments[0] = 0.05
    add_text(slide, layer_left + Inches(0.3), top2 + Inches(0.1), Inches(3), Inches(0.4),
             "② OpenClaw 平台中间层（重点）", size=15, bold=True, color=TEXT_LIGHT)
    add_text(slide, layer_left + Inches(0.3), top2 + Inches(0.55), Inches(2.8), Inches(0.4),
             "多 Agent 协同调度", size=11, color=NJU_PURPLE_LIGHT)
    cap = [("🧠 多模态理解", "截图 OCR · 二维码\n转账界面识别"),
           ("🔧 工具调用", "URL 黑名单 · 向量检索\n外部反诈接口 · 脱敏"),
           ("🔀 流程编排", "主控 Agent 调度\n案例回流自动化")]
    cw = Inches(2.9)
    cl = layer_left + Inches(3.5)
    for i, (k, v) in enumerate(cap):
        x = cl + (cw + Inches(0.1)) * i
        b = add_round_rect(slide, x, top2 + Inches(0.2), cw, Inches(1.15),
                           fill=NJU_PURPLE_DARK, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
        b.adjustments[0] = 0.1
        add_text(slide, x, top2 + Inches(0.3), cw, Inches(0.4),
                 k, size=14, bold=True, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
        add_text(slide, x, top2 + Inches(0.75), cw, Inches(0.55),
                 v, size=10, color=NJU_PURPLE_LIGHT, align=PP_ALIGN.CENTER)

    # 向下箭头
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                     Inches(6.4), top2 + layer_h2 + Inches(0.05),
                                     Inches(0.5), Inches(0.4))
    set_fill(arrow2, NJU_PURPLE)
    arrow2.line.fill.background()

    # Layer 3 知识库 + 输出
    top3 = top2 + layer_h2 + Inches(0.55)
    layer_h3 = Inches(1.35)
    l3 = add_round_rect(slide, layer_left, top3, layer_w, layer_h3,
                        fill=NJU_PURPLE_LIGHT, line=NJU_PURPLE, line_w=Pt(1.0))
    l3.adjustments[0] = 0.05
    add_text(slide, layer_left + Inches(0.3), top3 + Inches(0.1), Inches(3), Inches(0.4),
             "③ 知识库 + 三段输出", size=15, bold=True, color=NJU_PURPLE_DARK)
    add_text(slide, layer_left + Inches(0.3), top3 + Inches(0.55), Inches(3), Inches(0.4),
             "向量库 + 结构化案例库", size=11, color=TEXT_GRAY)
    outs = [("🔴🟡🟢", "三级风险标签", RED),
            ("📚", "RAG 相似案例引证", NJU_PURPLE),
            ("📌", "三步行动建议\n含 81686110 · 96110", GREEN)]
    ow = Inches(2.9)
    ol = layer_left + Inches(3.5)
    for i, (icon, txt, col) in enumerate(outs):
        x = ol + (ow + Inches(0.1)) * i
        b = add_round_rect(slide, x, top3 + Inches(0.15), ow, Inches(1.05),
                           fill=TEXT_LIGHT, line=col, line_w=Pt(0.75))
        b.adjustments[0] = 0.1
        add_text(slide, x, top3 + Inches(0.2), Inches(0.7), Inches(0.7),
                 icon, size=20, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, x + Inches(0.6), top3 + Inches(0.25), ow - Inches(0.6), Inches(0.8),
                 txt, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    add_footer(slide, "01:00", "1/3 · 设计目的与思路")

    # ------------------ 第 4 页 · OpenClaw 多 Agent 编排（重点1） ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, 4, "03  OpenClaw 平台应用 · 多 Agent 编排",
                    "重点 · 五个 Agent / Workflow 各承担一类任务")

    # 左侧 Agent 流程图
    diag_left = Inches(0.5)
    diag_top = Inches(1.55)
    diag_w = Inches(7.3)
    diag_bg = add_round_rect(slide, diag_left, diag_top, diag_w, Inches(5.3),
                             fill=BG_GRAY, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
    diag_bg.adjustments[0] = 0.03

    # 输入框
    in_box = add_round_rect(slide, diag_left + Inches(2.3), diag_top + Inches(0.3),
                            Inches(2.7), Inches(0.6),
                            fill=TEXT_LIGHT, line=TEXT_GRAY, line_w=Pt(0.5))
    in_box.adjustments[0] = 0.2
    add_text_in_shape(in_box, "用户输入（文本/截图/链接）",
                      size=12, bold=True, color=TEXT_DARK)

    # 向下到主控
    a1 = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                diag_left + Inches(3.4), diag_top + Inches(1.0),
                                Inches(0.4), Inches(0.35))
    set_fill(a1, NJU_PURPLE)
    a1.line.fill.background()

    # 主控 Agent
    mc = add_round_rect(slide, diag_left + Inches(2.0), diag_top + Inches(1.5),
                        Inches(3.3), Inches(0.7),
                        fill=NJU_PURPLE, line=NJU_PURPLE_DARK, line_w=Pt(1.0))
    mc.adjustments[0] = 0.25
    add_text_in_shape(mc, "🎯 主控 Agent\n识别输入类型 → 路由",
                      size=13, bold=True, color=TEXT_LIGHT)

    # 三个子 Agent
    sub_top = diag_top + Inches(2.5)
    sub_w = Inches(2.15)
    sub_h = Inches(1.2)
    sub_gap = Inches(0.15)
    sub_left_start = diag_left + Inches(0.25)
    subs = [
        ("📝 文本风险 Agent", "RAG 检索 + 规则匹配\n校园话术分类", NJU_PURPLE),
        ("🖼 多模态截图 Agent", "OCR + 二维码\n场景识别", NJU_PURPLE),
        ("🔗 链接电话 Agent", "URL 风险接口\n短链还原 / 号码归属", NJU_PURPLE),
    ]
    for i, (name, desc, col) in enumerate(subs):
        x = sub_left_start + (sub_w + sub_gap) * i
        b = add_round_rect(slide, x, sub_top, sub_w, sub_h,
                           fill=TEXT_LIGHT, line=col, line_w=Pt(1.0))
        b.adjustments[0] = 0.12
        add_text(slide, x, sub_top + Inches(0.1), sub_w, Inches(0.4),
                 name, size=12, bold=True, color=NJU_PURPLE_DARK, align=PP_ALIGN.CENTER)
        add_text(slide, x, sub_top + Inches(0.5), sub_w, Inches(0.7),
                 desc, size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # 向下汇聚
    agg_top = sub_top + sub_h + Inches(0.4)
    agg = add_round_rect(slide, diag_left + Inches(2.0), agg_top,
                         Inches(3.3), Inches(0.55),
                         fill=NJU_PURPLE_LIGHT, line=NJU_PURPLE, line_w=Pt(0.75))
    agg.adjustments[0] = 0.3
    add_text_in_shape(agg, "🧩 聚合输出工作流：风险标签 + 案例 + 三步建议",
                      size=11, bold=True, color=NJU_PURPLE_DARK)

    # 回流工作流
    wf_top = agg_top + Inches(0.8)
    wf = add_round_rect(slide, diag_left + Inches(0.5), wf_top,
                        Inches(6.3), Inches(0.55),
                        fill=TEXT_LIGHT, line=GREEN, line_w=Pt(0.75))
    wf.adjustments[0] = 0.3
    add_text_in_shape(wf, "♻️ 案例回流 Workflow：用户授权 → 脱敏 Agent → 写入知识库",
                      size=11, bold=True, color=GREEN)

    # 右侧 OpenClaw 能力映射
    map_left = Inches(8.1)
    map_top = Inches(1.55)
    map_w = Inches(4.7)
    map_h = Inches(5.3)
    map_bg = add_round_rect(slide, map_left, map_top, map_w, map_h,
                            fill=NJU_PURPLE_DARK, line=None)
    map_bg.adjustments[0] = 0.03
    add_text(slide, map_left + Inches(0.3), map_top + Inches(0.2), map_w - Inches(0.6), Inches(0.5),
             "▼ 各 Agent 调用的 OpenClaw 能力",
             size=14, bold=True, color=TEXT_LIGHT)

    mappings = [
        ("主控 Agent", "流程编排"),
        ("文本风险 Agent", "工具调用 (RAG) + 知识库"),
        ("多模态截图 Agent", "多模态理解 (OCR/QR)"),
        ("链接电话 Agent", "工具调用 (URL/号码接口)"),
        ("聚合输出工作流", "流程编排"),
        ("案例回流 Workflow", "流程编排 + 工具调用 (脱敏)"),
    ]
    item_top = map_top + Inches(0.85)
    for i, (agent, cap) in enumerate(mappings):
        y = item_top + Inches(0.72) * i
        # 圆点
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     map_left + Inches(0.3), y + Inches(0.18),
                                     Inches(0.18), Inches(0.18))
        set_fill(dot, NJU_PURPLE_LIGHT)
        dot.line.fill.background()
        add_text(slide, map_left + Inches(0.6), y, map_w - Inches(0.6), Inches(0.32),
                 agent, size=12, bold=True, color=TEXT_LIGHT)
        add_text(slide, map_left + Inches(0.6), y + Inches(0.3), map_w - Inches(0.6), Inches(0.32),
                 f"→ {cap}", size=11, color=NJU_PURPLE_LIGHT)

    add_footer(slide, "01:00", "2/3 · 平台应用与实现")

    # ------------------ 第 5 页 · 当前进展 · 双路召回 ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, 5, "04  当前进展 · 双路召回原型",
                    "本地 Demo v0.2.1 已跑通 · 决赛阶段迁移到 OpenClaw 平台")

    # 左：双路召回公式图
    fl = Inches(0.5)
    ft = Inches(1.55)
    fw = Inches(6.2)
    fh = Inches(5.3)
    fbg = add_round_rect(slide, fl, ft, fw, fh,
                         fill=BG_GRAY, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
    fbg.adjustments[0] = 0.03
    add_text(slide, fl + Inches(0.3), ft + Inches(0.2), fw - Inches(0.6), Inches(0.5),
             "▼ 双路融合检测引擎", size=15, bold=True, color=NJU_PURPLE_DARK)

    # 用户输入
    ui = add_round_rect(slide, fl + Inches(2.1), ft + Inches(0.9),
                        Inches(2.0), Inches(0.5),
                        fill=TEXT_LIGHT, line=TEXT_GRAY, line_w=Pt(0.5))
    ui.adjustments[0] = 0.2
    add_text_in_shape(ui, "用户输入", size=12, bold=True, color=TEXT_DARK)

    # 两条路
    rule = add_round_rect(slide, fl + Inches(0.3), ft + Inches(1.8),
                          Inches(2.6), Inches(1.0),
                          fill=TEXT_LIGHT, line=YELLOW, line_w=Pt(1.0))
    rule.adjustments[0] = 0.12
    add_text(slide, fl + Inches(0.4), ft + Inches(1.9), Inches(2.4), Inches(0.4),
             "📐 规则路", size=13, bold=True, color=YELLOW, align=PP_ALIGN.CENTER)
    add_text(slide, fl + Inches(0.4), ft + Inches(2.3), Inches(2.4), Inches(0.5),
             "keywords + regex 命中\n→ rule_score",
             size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    vec = add_round_rect(slide, fl + Inches(3.3), ft + Inches(1.8),
                         Inches(2.6), Inches(1.0),
                         fill=TEXT_LIGHT, line=NJU_PURPLE, line_w=Pt(1.0))
    vec.adjustments[0] = 0.12
    add_text(slide, fl + Inches(3.4), ft + Inches(1.9), Inches(2.4), Inches(0.4),
             "🧠 向量路", size=13, bold=True, color=NJU_PURPLE, align=PP_ALIGN.CENTER)
    add_text(slide, fl + Inches(3.4), ft + Inches(2.3), Inches(2.4), Inches(0.5),
             "BGE-zh 余弦相似度\n→ vector_sim",
             size=10, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # 公式
    formula = add_round_rect(slide, fl + Inches(0.6), ft + Inches(3.15),
                             Inches(5.0), Inches(0.6),
                             fill=NJU_PURPLE, line=None)
    formula.adjustments[0] = 0.3
    add_text_in_shape(formula,
                      "综合分 = rule_score + 2.5 × vector_sim",
                      size=14, bold=True, color=TEXT_LIGHT)

    # 风险分档
    risks = [("🔴 高风险", "双路确认", RED),
             ("🟡 中等", "单路强信号", YELLOW),
             ("🟢 低风险", "双路均弱", GREEN)]
    rt = ft + Inches(3.95)
    for i, (lvl, cond, col) in enumerate(risks):
        x = fl + Inches(0.3) + Inches(1.95) * i
        b = add_round_rect(slide, x, rt, Inches(1.85), Inches(1.1),
                           fill=TEXT_LIGHT, line=col, line_w=Pt(1.0))
        b.adjustments[0] = 0.12
        add_text(slide, x, rt + Inches(0.15), Inches(1.85), Inches(0.5),
                 lvl, size=14, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(slide, x, rt + Inches(0.65), Inches(1.85), Inches(0.4),
                 cond, size=11, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    # 右：当前进展数据 + 决赛迁移
    rl = Inches(7.0)
    rw = Inches(5.8)
    rt2 = Inches(1.55)
    rh = Inches(5.3)
    rbg = add_round_rect(slide, rl, rt2, rw, rh,
                         fill=TEXT_LIGHT, line=NJU_PURPLE, line_w=Pt(1.0))
    rbg.adjustments[0] = 0.03
    add_text(slide, rl + Inches(0.3), rt2 + Inches(0.2), rw - Inches(0.6), Inches(0.5),
             "▼ 当前进展（本地 Demo v0.2.1）",
             size=15, bold=True, color=NJU_PURPLE_DARK)

    stats = [
        ("📚 50 类知识库", "国家反诈 10 类 + 新型 5 类 + 校园特化 12 类"),
        ("🎯 8 / 8 冒烟测试", "TOP-1 全部命中正确 case"),
        ("⚡ ~180ms 单次延迟", "CPU 推理 · 嵌入缓存毫秒级二次启动"),
        ("🔄 降级链", "BGE-zh → TF-IDF → 仅规则"),
        ("📂 来源可追溯", "每条 case 强制 source 字段引用公开材料"),
        ("🛡 CI 已绿", "GitHub Actions · MIT + CC BY 4.0 双许可"),
    ]
    for i, (k, v) in enumerate(stats):
        y = rt2 + Inches(0.85) + Inches(0.55) * i
        add_text(slide, rl + Inches(0.3), y, rw - Inches(0.6), Inches(0.3),
                 k, size=12, bold=True, color=NJU_PURPLE)
        add_text(slide, rl + Inches(0.3), y + Inches(0.25), rw - Inches(0.6), Inches(0.3),
                 v, size=10, color=TEXT_GRAY)

    # 底部迁移条
    mig = add_round_rect(slide, rl + Inches(0.2), rt2 + rh - Inches(0.8),
                         rw - Inches(0.4), Inches(0.6),
                         fill=NJU_PURPLE_LIGHT, line=NJU_PURPLE, line_w=Pt(0.5))
    mig.adjustments[0] = 0.3
    add_text_in_shape(mig,
                      "▶ 决赛迁移：本地规则/TF-IDF → OpenClaw 原生 RAG + 多 Agent",
                      size=12, bold=True, color=NJU_PURPLE_DARK)

    add_footer(slide, "00:40", "2/3 · 平台应用与实现")

    # ------------------ 第 6 页 · 现场 Demo · token 环节 ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, 6, "05  现场 Demo · token 使用环节（重点）",
                    "屏幕共享 OpenClaw 平台 → 实操检测 → 平台调度日志可见")

    # 左侧：演示输入卡片
    ll = Inches(0.5)
    lt = Inches(1.55)
    lw = Inches(5.6)
    lh = Inches(5.3)
    lbg = add_round_rect(slide, ll, lt, lw, lh,
                         fill=BG_GRAY, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
    lbg.adjustments[0] = 0.03

    # 输入框 mockup
    add_text(slide, ll + Inches(0.3), lt + Inches(0.2), lw - Inches(0.6), Inches(0.4),
             "▼ 现场演示输入", size=14, bold=True, color=NJU_PURPLE_DARK)
    inp = add_round_rect(slide, ll + Inches(0.3), lt + Inches(0.7),
                         lw - Inches(0.6), Inches(0.95),
                         fill=TEXT_LIGHT, line=TEXT_GRAY, line_w=Pt(0.75))
    inp.adjustments[0] = 0.08
    add_text_in_shape(inp,
                      "「辅导员让我代收班费，要我先\n转 800 给他他周一返我」",
                      size=13, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # 输出卡片 mockup（红色高风险）
    add_text(slide, ll + Inches(0.3), lt + Inches(1.85), lw - Inches(0.6), Inches(0.4),
             "▼ 预期输出（OpenClaw 端到端）", size=14, bold=True, color=NJU_PURPLE_DARK)

    # 红色风险条
    risk_bar = add_round_rect(slide, ll + Inches(0.3), lt + Inches(2.35),
                              lw - Inches(0.6), Inches(0.55),
                              fill=RED, line=None)
    risk_bar.adjustments[0] = 0.2
    add_text_in_shape(risk_bar,
                      "🔴 高风险 · 命中 KB-008 冒充辅导员 · 综合分 4.2",
                      size=12, bold=True, color=TEXT_LIGHT)

    # 三步建议
    advices = [
        "1. 立即停止转账，通过电话/当面与辅导员核实",
        "2. 八个凡是·第 6 条：自称领导先嘘寒问暖再要转账",
        "3. 已转账 → 96110 / 南大保卫处 81686110 / 报警 110",
    ]
    for i, t in enumerate(advices):
        y = lt + Inches(3.1) + Inches(0.55) * i
        b = add_round_rect(slide, ll + Inches(0.3), y, lw - Inches(0.6), Inches(0.5),
                           fill=TEXT_LIGHT, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
        b.adjustments[0] = 0.15
        add_text(slide, ll + Inches(0.45), y + Inches(0.1), lw - Inches(0.8), Inches(0.35),
                 t, size=11, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # 兜底说明
    bk = add_round_rect(slide, ll + Inches(0.3), lt + lh - Inches(0.6),
                        lw - Inches(0.6), Inches(0.4),
                        fill=YELLOW, line=None)
    bk.adjustments[0] = 0.2
    add_text_in_shape(bk, "⚠ 兜底：网络异常立即切本地 Gradio 或播 Demo.mov",
                      size=10, bold=True, color=TEXT_DARK)

    # 右侧：5 步 token 流转脚本
    rl = Inches(6.4)
    rt2 = Inches(1.55)
    rw = Inches(6.4)
    rh = Inches(5.3)
    rbg = add_round_rect(slide, rl, rt2, rw, rh,
                         fill=NJU_PURPLE_DARK, line=None)
    rbg.adjustments[0] = 0.03
    add_text(slide, rl + Inches(0.3), rt2 + Inches(0.2), rw - Inches(0.6), Inches(0.5),
             "▼ OpenClaw 平台 token 流转 · 5 步演示脚本",
             size=14, bold=True, color=TEXT_LIGHT)

    steps = [
        ("Step 1", "粘贴输入到 OpenClaw Web 控制台", "🪙 输入 token 解析"),
        ("Step 2", "主控 Agent 识别 → 路由至文本风险 Agent",
         "🪙 流程编排 token · 平台日志可见"),
        ("Step 3", "文本风险 Agent 调用 RAG：BGE-zh 向量召回",
         "🪙 嵌入 token + 检索 token"),
        ("Step 4", "工具调用：URL/号码黑名单接口（本例跳过）",
         "🪙 工具调用 token"),
        ("Step 5", "聚合输出 🔴高风险卡片 · 全链路 token 消耗可见",
         "🪙 推理 token · 计费透明"),
    ]
    for i, (no, desc, tok) in enumerate(steps):
        y = rt2 + Inches(0.8) + Inches(0.9) * i
        # 步骤号圆
        nob = add_round_rect(slide, rl + Inches(0.3), y, Inches(0.8), Inches(0.7),
                             fill=NJU_PURPLE, line=NJU_PURPLE_LIGHT, line_w=Pt(1.0))
        nob.adjustments[0] = 0.3
        add_text_in_shape(nob, no, size=12, bold=True, color=TEXT_LIGHT)
        # 说明
        add_text(slide, rl + Inches(1.2), y, rw - Inches(1.5), Inches(0.35),
                 desc, size=11, bold=True, color=TEXT_LIGHT)
        add_text(slide, rl + Inches(1.2), y + Inches(0.32), rw - Inches(1.5), Inches(0.35),
                 tok, size=10, color=NJU_PURPLE_LIGHT)

    add_footer(slide, "01:10", "3/3 · 未来效果与现场演示")

    # ------------------ 第 7 页 · 决赛交付 + 差异 + 隐私 ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_page_header(slide, 7, "06  决赛交付 · 差异化定位 · 隐私边界",
                    "六项决赛交付物 · 与国家反诈 App 的互补关系")

    # 左：决赛交付 6 项
    ll = Inches(0.5)
    lt = Inches(1.55)
    lw = Inches(6.0)
    lh = Inches(5.3)
    lbg = add_round_rect(slide, ll, lt, lw, lh,
                         fill=BG_GRAY, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
    lbg.adjustments[0] = 0.03
    add_text(slide, ll + Inches(0.3), lt + Inches(0.2), lw - Inches(0.6), Inches(0.4),
             "▼ 决赛阶段交付清单", size=15, bold=True, color=NJU_PURPLE_DARK)

    deliverables = [
        ("🚀", "OpenClaw 平台多 Agent 端到端演示 + 微信小程序前端"),
        ("📚", "知识库扩充至 200+ 条 · 含用户回流脱敏工作流"),
        ("🛡", "数据隐私与 AI 安全说明书"),
        ("🌐", "代码 / 提示词 / 知识库 / 文档全量开源"),
        ("🎬", "3 分钟路演视频 + 典型案例对话演示"),
        ("🤝", "联系学工部、保卫处寻求真实案例支持"),
    ]
    for i, (icon, t) in enumerate(deliverables):
        y = lt + Inches(0.8) + Inches(0.72) * i
        b = add_round_rect(slide, ll + Inches(0.3), y, lw - Inches(0.6), Inches(0.62),
                           fill=TEXT_LIGHT, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
        b.adjustments[0] = 0.15
        add_text(slide, ll + Inches(0.45), y + Inches(0.1), Inches(0.55), Inches(0.45),
                 icon, size=18, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, ll + Inches(1.05), y + Inches(0.13), lw - Inches(1.4), Inches(0.4),
                 t, size=12, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # 右：差异表 + 隐私边界
    rl = Inches(6.8)
    rt2 = Inches(1.55)
    rw = Inches(6.0)

    # 差异
    dbg = add_round_rect(slide, rl, rt2, rw, Inches(2.6),
                         fill=NJU_PURPLE_LIGHT, line=NJU_PURPLE, line_w=Pt(0.5))
    dbg.adjustments[0] = 0.05
    add_text(slide, rl + Inches(0.3), rt2 + Inches(0.15), rw - Inches(0.6), Inches(0.4),
             "▼ 与国家反诈中心 App 差异", size=14, bold=True, color=NJU_PURPLE_DARK)
    # 简易两栏对比
    diffs = [
        ("通用 vs 校园语境", "辅导员/教务/奖助/二手/租房/老乡群高仿号"),
        ("通用紧急电话", "本地：南大保卫处 81686110 直拨"),
        ("独立 App", "微信小程序 + 嵌入南大官方 App 扩展能力"),
    ]
    for i, (l, r) in enumerate(diffs):
        y = rt2 + Inches(0.65) + Inches(0.6) * i
        add_text(slide, rl + Inches(0.4), y, Inches(2.0), Inches(0.35),
                 f"📌 {l}", size=11, bold=True, color=NJU_PURPLE_DARK)
        add_text(slide, rl + Inches(2.5), y, rw - Inches(2.8), Inches(0.55),
                 r, size=10, color=TEXT_DARK)

    # 隐私边界
    pt = rt2 + Inches(2.8)
    pbg = add_round_rect(slide, rl, pt, rw, Inches(2.5),
                         fill=NJU_PURPLE_DARK, line=None)
    pbg.adjustments[0] = 0.05
    add_text(slide, rl + Inches(0.3), pt + Inches(0.15), rw - Inches(0.6), Inches(0.4),
             "▼ 数据与隐私安全（AI 安全之星奖关注）",
             size=14, bold=True, color=TEXT_LIGHT)
    privacy = [
        "✓ 即用即弃：默认不持久化用户输入",
        "✓ 用户授权方启动回流；脱敏 Agent 自动替换 PII",
        "✓ 案例 ID + 匹配依据可追溯",
        "✓ 仅识别与提示，不替代决策、不拦截通讯",
    ]
    for i, t in enumerate(privacy):
        y = pt + Inches(0.65) + Inches(0.4) * i
        add_text(slide, rl + Inches(0.4), y, rw - Inches(0.6), Inches(0.35),
                 t, size=11, color=NJU_PURPLE_LIGHT)

    add_footer(slide, "00:30", "3/3 · 未来效果与现场演示")

    # ------------------ 第 8 页 · 团队 + 致谢 + Q&A ------------------
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NJU_PURPLE_DARK)
    # 大标题
    add_text(slide, Inches(0.7), Inches(0.6), Inches(12), Inches(0.5),
             "07  Q & A · 期待评委质询", size=14, color=NJU_PURPLE_LIGHT, bold=True)
    add_text(slide, Inches(0.7), Inches(1.1), Inches(12), Inches(1.0),
             "感谢评审老师", size=48, color=TEXT_LIGHT, bold=True)
    add_text(slide, Inches(0.7), Inches(2.2), Inches(12), Inches(0.6),
             "南大数智安全官 NJU Guardian · 期待您的指导与建议",
             size=18, color=NJU_PURPLE_LIGHT)

    # 三栏底部信息
    cols = [
        ("👤 团队", ["作者（***REMOVED***）· 个人参赛",
                    "南京大学 · 经济学（拔尖计划）",
                    "联系：***REMOVED***"]),
        ("📚 数据来源", ["国家反诈中心《2023 版宣传手册》",
                       "公安部刑侦局公开预警通报",
                       "央视 · 新华社 · 高校公开通报",
                       "南大保卫处反诈材料"]),
        ("📞 紧急联系", ["96110 · 反诈劝阻专线",
                       "南大保卫处 · 81686110",
                       "12381 · 涉诈短信预警",
                       "110 · 报警"]),
    ]
    cw = Inches(4.0)
    cl = Inches(0.5)
    ct = Inches(3.6)
    ch = Inches(3.0)
    for i, (title, lines) in enumerate(cols):
        x = cl + (cw + Inches(0.15)) * i
        b = add_round_rect(slide, x, ct, cw, ch,
                           fill=NJU_PURPLE, line=NJU_PURPLE_LIGHT, line_w=Pt(0.5))
        b.adjustments[0] = 0.05
        add_text(slide, x + Inches(0.3), ct + Inches(0.2), cw - Inches(0.6), Inches(0.45),
                 title, size=15, bold=True, color=TEXT_LIGHT)
        for j, line in enumerate(lines):
            y = ct + Inches(0.8) + Inches(0.45) * j
            add_text(slide, x + Inches(0.3), y, cw - Inches(0.6), Inches(0.4),
                     line, size=12, color=NJU_PURPLE_LIGHT)

    # 仓库地址条
    rep = add_round_rect(slide, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.5),
                         fill=NJU_PURPLE_LIGHT, line=None)
    rep.adjustments[0] = 0.3
    add_text_in_shape(rep,
                      "🌐 开源仓库：github.com/hjjbh1314/nju-guardian   ·   "
                      "MIT (代码) + CC BY 4.0 (知识库与视觉素材)",
                      size=12, bold=True, color=NJU_PURPLE_DARK)

    prs.save(output_path)
    print(f"✅ PPT 已生成：{output_path}")


if __name__ == "__main__":
    import os
    out = "/Users/haiwenbao/Documents/OpenClaw应用创新大赛/答辩PPT_南大数智安全官_v1.pptx"
    build_deck(out)
    sz = os.path.getsize(out) / 1024
    print(f"   文件大小：{sz:.1f} KB")
